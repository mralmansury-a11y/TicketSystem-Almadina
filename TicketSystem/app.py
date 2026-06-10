from flask import Flask, render_template, request, jsonify, session, send_file
from datetime import datetime
import sqlite3, io, os

app = Flask(__name__)
app.secret_key = "amp_support_secret_2024"
DB_PATH = os.path.join(os.path.dirname(__file__), "support.db")
TICKET_PREFIX = "T-"


# ══════════════════════════════════════════════
#  DB
# ══════════════════════════════════════════════
def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    conn = get_db(); c = conn.cursor()
    c.executescript("""
    CREATE TABLE IF NOT EXISTS tickets (
        id           INTEGER PRIMARY KEY AUTOINCREMENT,
        created_at   TEXT NOT NULL,
        name         TEXT NOT NULL,
        dept         TEXT NOT NULL,
        issue        TEXT NOT NULL,
        description  TEXT,
        ticket_no    TEXT UNIQUE NOT NULL,
        status       TEXT DEFAULT 'جديد',
        completed_at TEXT,
        tech         TEXT DEFAULT '',
        updated_at   TEXT,
        email        TEXT,
        phone        TEXT,
        priority     TEXT DEFAULT 'عادي',
        solution     TEXT DEFAULT ''
    );
    CREATE TABLE IF NOT EXISTS users (
        id       INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL,
        password TEXT NOT NULL,
        role     TEXT NOT NULL,
        fullname TEXT,
        active   INTEGER DEFAULT 1
    );
    CREATE TABLE IF NOT EXISTS departments (
        id   INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT UNIQUE NOT NULL,
        active INTEGER DEFAULT 1
    );
    CREATE TABLE IF NOT EXISTS issue_types (
        id   INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT UNIQUE NOT NULL,
        active INTEGER DEFAULT 1
    );
    CREATE TABLE IF NOT EXISTS ticket_notes (
        id         INTEGER PRIMARY KEY AUTOINCREMENT,
        ticket_no  TEXT NOT NULL,
        author     TEXT NOT NULL,
        note       TEXT NOT NULL,
        created_at TEXT NOT NULL
    );
    """)

    # Add solution column if it doesn't exist (migration for existing DBs)
    try:
        c.execute("ALTER TABLE tickets ADD COLUMN solution TEXT DEFAULT ''")
        conn.commit()
    except Exception:
        pass

    c.execute("SELECT COUNT(*) FROM users")
    if c.fetchone()[0] == 0:
        c.executemany("INSERT INTO users (username,password,role,fullname) VALUES (?,?,?,?)", [
            ("admin",  "admin123", "admin", "مدير النظام"),
            ("tech1",  "tech123",  "tech",  "أحمد الفني"),
            ("tech2",  "tech456",  "tech",  "محمد الدعم"),
            ("tech3",  "tech789",  "tech",  "علي التقني"),
        ])

    c.execute("SELECT COUNT(*) FROM departments")
    if c.fetchone()[0] == 0:
        for d in ["قسم تقنية المعلومات","قسم الموارد البشرية","قسم المحاسبة","قسم الإدارة","قسم المشاريع","قسم المبيعات","قسم القانوني","قسم العمليات","قسم التسويق"]:
            c.execute("INSERT INTO departments (name) VALUES (?)", (d,))

    c.execute("SELECT COUNT(*) FROM issue_types")
    if c.fetchone()[0] == 0:
        for i in ["عطل شبكة","عطل طابعة","مشكلة برنامج","عطل حاسوب","مشكلة بريد إلكتروني","طلب صيانة","مشكلة أخرى","إعادة تثبيت نظام","مشكلة شبكة WiFi"]:
            c.execute("INSERT INTO issue_types (name) VALUES (?)", (i,))

    conn.commit(); conn.close()


def now_str():
    return datetime.now().strftime("%Y-%m-%d %H:%M")


def generate_ticket_number():
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT ticket_no FROM tickets ORDER BY id DESC LIMIT 1")
    row = c.fetchone(); conn.close()
    if not row: return f"{TICKET_PREFIX}001"
    try: num = int(row["ticket_no"].replace(TICKET_PREFIX,"")) + 1
    except: num = 1
    return f"{TICKET_PREFIX}{str(num).zfill(3)}"


def require_admin():
    return session.get("role") == "admin"

# ══════════════════════════════════════════════
#  Pages
# ══════════════════════════════════════════════
@app.route("/")
def index():
    return render_template("index.html")


# ══════════════════════════════════════════════
#  Auth
# ══════════════════════════════════════════════
@app.route("/api/login", methods=["POST"])
def api_login():
    d = request.json
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT * FROM users WHERE LOWER(username)=? AND password=? AND active=1",
              (d.get("username","").strip().lower(), d.get("password","").strip()))
    u = c.fetchone(); conn.close()
    if u:
        session["user"] = u["username"]; session["role"] = u["role"]
        return jsonify({"ok":True,"user":u["username"],"role":u["role"],"fullname":u["fullname"] or u["username"]})
    return jsonify({"ok":False})

@app.route("/api/logout", methods=["POST"])
def api_logout():
    session.clear(); return jsonify({"ok":True})


# ══════════════════════════════════════════════
#  Lookups
# ══════════════════════════════════════════════
@app.route("/api/lookups")
def api_lookups():
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT name FROM departments WHERE active=1 ORDER BY name")
    depts = [r["name"] for r in c.fetchall()]
    c.execute("SELECT name FROM issue_types WHERE active=1 ORDER BY name")
    issues = [r["name"] for r in c.fetchall()]
    conn.close()
    return jsonify({"depts":depts,"issues":issues})

@app.route("/api/techs")
def api_techs():
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT username, fullname FROM users WHERE role='tech' AND active=1")
    rows = c.fetchall(); conn.close()
    return jsonify([{"username":r["username"],"fullname":r["fullname"] or r["username"]} for r in rows])


# ══════════════════════════════════════════════
#  User Management (admin only)
# ══════════════════════════════════════════════
@app.route("/api/users")
def api_users():
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT id,username,role,fullname,active FROM users ORDER BY id")
    rows = [dict(r) for r in c.fetchall()]; conn.close()
    return jsonify(rows)

@app.route("/api/users", methods=["POST"])
def api_add_user():
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    d = request.json
    if not d.get("username") or not d.get("password") or not d.get("role"):
        return jsonify({"success":False,"error":"بيانات ناقصة"}), 400
    conn = get_db(); c = conn.cursor()
    try:
        c.execute("INSERT INTO users (username,password,role,fullname) VALUES (?,?,?,?)",
                  (d["username"].strip().lower(), d["password"], d["role"], d.get("fullname","").strip()))
        conn.commit(); conn.close()
        return jsonify({"success":True})
    except sqlite3.IntegrityError:
        conn.close(); return jsonify({"success":False,"error":"اسم المستخدم موجود مسبقاً"})

@app.route("/api/users/<int:uid>", methods=["PUT"])
def api_update_user(uid):
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    d = request.json; conn = get_db(); c = conn.cursor()
    fields, vals = [], []
    if "fullname" in d: fields.append("fullname=?"); vals.append(d["fullname"])
    if "role"     in d: fields.append("role=?");     vals.append(d["role"])
    if "password" in d and d["password"]: fields.append("password=?"); vals.append(d["password"])
    if "active"   in d: fields.append("active=?");   vals.append(int(d["active"]))
    if fields:
        vals.append(uid)
        c.execute(f"UPDATE users SET {','.join(fields)} WHERE id=?", vals)
        conn.commit()
    conn.close(); return jsonify({"success":True})

@app.route("/api/users/<int:uid>", methods=["DELETE"])
def api_delete_user(uid):
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    conn = get_db(); c = conn.cursor()
    c.execute("UPDATE users SET active=0 WHERE id=?", (uid,))
    conn.commit(); conn.close(); return jsonify({"success":True})


# ══════════════════════════════════════════════
#  Department Management (admin)
# ══════════════════════════════════════════════
@app.route("/api/departments")
def api_depts():
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT id,name,active FROM departments ORDER BY name")
    rows = [dict(r) for r in c.fetchall()]; conn.close()
    return jsonify(rows)

@app.route("/api/departments", methods=["POST"])
def api_add_dept():
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    name = (request.json or {}).get("name","").strip()
    if not name: return jsonify({"success":False,"error":"الاسم مطلوب"})
    conn = get_db(); c = conn.cursor()
    try:
        c.execute("INSERT INTO departments (name) VALUES (?)", (name,))
        conn.commit(); conn.close(); return jsonify({"success":True})
    except sqlite3.IntegrityError:
        conn.close(); return jsonify({"success":False,"error":"القسم موجود مسبقاً"})

@app.route("/api/departments/<int:did>", methods=["PUT"])
def api_update_dept(did):
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    d = request.json; conn = get_db(); c = conn.cursor()
    if "name" in d:   c.execute("UPDATE departments SET name=?   WHERE id=?", (d["name"], did))
    if "active" in d: c.execute("UPDATE departments SET active=? WHERE id=?", (int(d["active"]), did))
    conn.commit(); conn.close(); return jsonify({"success":True})

@app.route("/api/departments/<int:did>", methods=["DELETE"])
def api_delete_dept(did):
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    conn = get_db(); c = conn.cursor()
    c.execute("UPDATE departments SET active=0 WHERE id=?", (did,))
    conn.commit(); conn.close(); return jsonify({"success":True})


# ══════════════════════════════════════════════
#  Issue Type Management (admin)
# ══════════════════════════════════════════════
@app.route("/api/issue-types")
def api_issues():
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT id,name,active FROM issue_types ORDER BY name")
    rows = [dict(r) for r in c.fetchall()]; conn.close()
    return jsonify(rows)

@app.route("/api/issue-types", methods=["POST"])
def api_add_issue():
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    name = (request.json or {}).get("name","").strip()
    if not name: return jsonify({"success":False,"error":"الاسم مطلوب"})
    conn = get_db(); c = conn.cursor()
    try:
        c.execute("INSERT INTO issue_types (name) VALUES (?)", (name,))
        conn.commit(); conn.close(); return jsonify({"success":True})
    except sqlite3.IntegrityError:
        conn.close(); return jsonify({"success":False,"error":"النوع موجود مسبقاً"})

@app.route("/api/issue-types/<int:iid>", methods=["PUT"])
def api_update_issue(iid):
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    d = request.json; conn = get_db(); c = conn.cursor()
    if "name"   in d: c.execute("UPDATE issue_types SET name=?   WHERE id=?", (d["name"], iid))
    if "active" in d: c.execute("UPDATE issue_types SET active=? WHERE id=?", (int(d["active"]), iid))
    conn.commit(); conn.close(); return jsonify({"success":True})

@app.route("/api/issue-types/<int:iid>", methods=["DELETE"])
def api_delete_issue(iid):
    if not require_admin(): return jsonify({"error":"unauthorized"}), 403
    conn = get_db(); c = conn.cursor()
    c.execute("UPDATE issue_types SET active=0 WHERE id=?", (iid,))
    conn.commit(); conn.close(); return jsonify({"success":True})


# ══════════════════════════════════════════════
#  Tickets
# ══════════════════════════════════════════════
@app.route("/api/tickets", methods=["POST"])
def api_add_ticket():
    d = request.json; ticket_no = generate_ticket_number(); ts = now_str()
    conn = get_db(); c = conn.cursor()
    c.execute("""INSERT INTO tickets
       (created_at,name,dept,issue,description,ticket_no,status,tech,updated_at,email,phone,priority)
       VALUES (?,?,?,?,?,?,?,?,?,?,?,?)""",
       (ts,d["name"],d["dept"],d["issue"],d.get("desc",""),ticket_no,"جديد","",ts,
        d.get("email",""),d.get("phone",""),d.get("priority","عادي")))
    conn.commit(); conn.close()
    return jsonify({"success":True,"ticket":ticket_no})

@app.route("/api/tickets")
def api_get_tickets():
    conn = get_db(); c = conn.cursor()
    sf = request.args.get("status"); tf = request.args.get("tech")
    if sf:   c.execute("SELECT * FROM tickets WHERE status=? ORDER BY id DESC",(sf,))
    elif tf: c.execute("SELECT * FROM tickets WHERE tech=? ORDER BY id DESC",(tf,))
    else:    c.execute("SELECT * FROM tickets ORDER BY id DESC")
    rows = [dict(r) for r in c.fetchall()]; conn.close()
    return jsonify(rows)

@app.route("/api/tickets/<ticket_no>")
def api_get_ticket(ticket_no):
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT * FROM tickets WHERE ticket_no=?",(ticket_no,))
    row = c.fetchone()
    c.execute("SELECT * FROM ticket_notes WHERE ticket_no=? ORDER BY id ASC",(ticket_no,))
    notes = [dict(n) for n in c.fetchall()]; conn.close()
    if not row: return jsonify({"error":"not found"}),404
    result = dict(row); result["notes"] = notes
    return jsonify(result)

@app.route("/api/tickets/<ticket_no>/assign", methods=["POST"])
def api_assign(ticket_no):
    tech_name = request.json.get("tech")
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT status FROM tickets WHERE ticket_no=?",(ticket_no,))
    row = c.fetchone()
    new_status = row["status"] if row and row["status"]!="جديد" else "قيد التنفيذ"
    c.execute("UPDATE tickets SET tech=?,status=?,updated_at=? WHERE ticket_no=?",
              (tech_name,new_status,now_str(),ticket_no))
    conn.commit(); conn.close(); return jsonify({"success":True})

@app.route("/api/tickets/<ticket_no>/status", methods=["POST"])
def api_update_status(ticket_no):
    d = request.json
    status = d.get("status")
    solution = d.get("solution", "")
    conn = get_db(); c = conn.cursor(); ts = now_str()
    if status == "مكتمل":
        c.execute("UPDATE tickets SET status=?,completed_at=?,updated_at=?,solution=? WHERE ticket_no=?",
                  (status, ts, ts, solution, ticket_no))
    else:
        c.execute("UPDATE tickets SET status=?,updated_at=? WHERE ticket_no=?",(status,ts,ticket_no))
    conn.commit(); conn.close(); return jsonify({"success":True})

@app.route("/api/tickets/<ticket_no>/notes", methods=["POST"])
def api_add_note(ticket_no):
    d = request.json; conn = get_db(); c = conn.cursor()
    c.execute("INSERT INTO ticket_notes (ticket_no,author,note,created_at) VALUES (?,?,?,?)",
              (ticket_no,d.get("author",""),d.get("note",""),now_str()))
    conn.commit(); conn.close(); return jsonify({"success":True})


# ══════════════════════════════════════════════
#  Stats
# ══════════════════════════════════════════════
@app.route("/api/stats/quick")
def api_quick_stats():
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT status,COUNT(*) as cnt FROM tickets GROUP BY status")
    stats = {"done":0,"process":0,"new":0}
    for r in c.fetchall():
        if r["status"]=="مكتمل":       stats["done"]    = r["cnt"]
        elif r["status"]=="قيد التنفيذ": stats["process"] = r["cnt"]
        elif r["status"]=="جديد":       stats["new"]     = r["cnt"]
    conn.close(); return jsonify(stats)

@app.route("/api/stats/advanced")
def api_advanced_stats():
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT * FROM tickets"); rows = [dict(r) for r in c.fetchall()]; conn.close()
    total     = len(rows)
    completed = sum(1 for r in rows if r["status"]=="مكتمل")
    in_prog   = sum(1 for r in rows if r["status"]=="قيد التنفيذ")
    new       = sum(1 for r in rows if r["status"]=="جديد")
    dept_map  = {}
    for r in rows: dept_map[r["dept"]]  = dept_map.get(r["dept"],0)+1
    issue_map = {}
    for r in rows: issue_map[r["issue"]] = issue_map.get(r["issue"],0)+1
    tech_map  = {}
    for r in rows:
        if r.get("tech"): tech_map[r["tech"]] = tech_map.get(r["tech"],0)+1
    daily_map = {}
    for r in rows:
        day = str(r["created_at"])[:10]; daily_map[day] = daily_map.get(day,0)+1
    times = []
    for r in rows:
        if r["status"]=="مكتمل" and r["created_at"] and r["completed_at"]:
            try:
                t1 = datetime.strptime(r["created_at"][:16],"%Y-%m-%d %H:%M")
                t2 = datetime.strptime(r["completed_at"][:16],"%Y-%m-%d %H:%M")
                times.append((t2-t1).total_seconds()/3600)
            except: pass
    ct = {"avg":round(sum(times)/len(times),1) if times else 0,
          "min":round(min(times),1) if times else 0,
          "max":round(max(times),1) if times else 0}
    return jsonify({"total":total,"completed":completed,"inProgress":in_prog,"new":new,
        "deptChart":[{"dept":k,"count":v} for k,v in dept_map.items()],
        "issueChart":[{"issue":k,"count":v} for k,v in issue_map.items()],
        "techChart":[{"tech":k,"count":v} for k,v in tech_map.items()],
        "dailyChart":[{"date":k,"count":v} for k,v in sorted(daily_map.items())],
        "completionTime":ct})

@app.route("/api/activities")
def api_activities():
    conn = get_db(); c = conn.cursor()
    c.execute("SELECT * FROM tickets ORDER BY id DESC LIMIT 15")
    rows = c.fetchall(); conn.close()
    result = []
    for r in rows:
        action = "تم الإكمال" if r["status"]=="مكتمل" else "بلاغ جديد" if r["status"]=="جديد" else "قيد المعالجة"
        result.append({"ticket":r["ticket_no"],"action":action,"time":r["updated_at"] or r["created_at"],
                       "user":r["name"],"tech":r["tech"] or "","status":r["status"]})
    return jsonify(result)


# ══════════════════════════════════════════════
#  EXPORT: Excel
# ══════════════════════════════════════════════
@app.route("/api/export/excel")
def api_export_excel():
    if not require_admin(): return jsonify({"error":"unauthorized"}),403
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    conn = get_db(); c = conn.cursor()
    c.execute("SELECT * FROM tickets ORDER BY id")
    tickets = [dict(r) for r in c.fetchall()]
    c.execute("SELECT tn.*, t.name as employee FROM ticket_notes tn LEFT JOIN tickets t ON tn.ticket_no=t.ticket_no ORDER BY tn.id")
    notes = [dict(r) for r in c.fetchall()]
    c.execute("SELECT * FROM users ORDER BY role,username")
    users = [dict(r) for r in c.fetchall()]
    conn.close()

    wb = Workbook()

    # ── Palette ──
    DARK_BLUE = "1E3A5F"; MID_BLUE = "2563EB"; LIGHT_BLUE = "DBEAFE"
    GREEN  = "059652"; LIGHT_GREEN = "D1FAE5"
    ORANGE = "D97706"; LIGHT_ORANGE= "FEF3C7"
    RED    = "BE123C"; LIGHT_RED   = "FFE4E6"
    WHITE  = "FFFFFF"; GRAY_ROW    = "F8FAFC"; GRAY_HDR   = "1E293B"
    TEXT_DARK = "111827"; TEXT_MID = "374151"
    PURPLE = "7C3AED"; LIGHT_PURPLE = "EDE9FE"

    thin = Side(style="thin",color="CBD5E1")
    border_all = Border(left=thin,right=thin,top=thin,bottom=thin)

    def hdr_font(sz=11, bold=True, color=WHITE):
        return Font(name="Arial",size=sz,bold=bold,color=color)
    def cell_font(sz=10, bold=False, color=TEXT_DARK):
        return Font(name="Arial",size=sz,bold=bold,color=color)
    def fill(hex_color):
        return PatternFill("solid", start_color=hex_color, fgColor=hex_color)
    def center():
        return Alignment(horizontal="center",vertical="center",wrap_text=True,readingOrder=2)
    def right():
        return Alignment(horizontal="right",vertical="center",wrap_text=True,readingOrder=2)

    def write_header(ws, headers, col_widths, bg=DARK_BLUE):
        ws.row_dimensions[1].height = 32
        for col,(hdr,w) in enumerate(zip(headers,col_widths),1):
            cell = ws.cell(row=1,column=col,value=hdr)
            cell.font = hdr_font(); cell.fill = fill(bg)
            cell.alignment = center(); cell.border = border_all
            ws.column_dimensions[get_column_letter(col)].width = w

    STATUS_COLOR = {"جديد":RED,"قيد التنفيذ":ORANGE,"مكتمل":GREEN}
    STATUS_BG    = {"جديد":LIGHT_RED,"قيد التنفيذ":LIGHT_ORANGE,"مكتمل":LIGHT_GREEN}
    PRI_COLOR    = {"عادي":MID_BLUE,"عاجل":ORANGE,"حرج":RED}

    # ── Sheet 1: Tickets ──
    ws1 = wb.active; ws1.title = "البلاغات"
    ws1.sheet_view.rightToLeft = True
    headers1 = ["رقم البلاغ","الاسم","القسم","نوع المشكلة","الأولوية","الحالة","الفني المعيّن","البريد الإلكتروني","الهاتف","تاريخ الإنشاء","تاريخ الإكمال","وصف المشكلة","الحل المُطبَّق"]
    widths1   = [12,18,20,22,10,14,16,22,16,18,18,40,40]
    write_header(ws1, headers1, widths1)

    for i,t in enumerate(tickets,2):
        ws1.row_dimensions[i].height = 22
        status = t.get("status","جديد")
        pri    = t.get("priority","عادي")
        row_bg = GRAY_ROW if i%2==0 else WHITE
        vals = [t["ticket_no"],t["name"],t["dept"],t["issue"],pri,status,
                t.get("tech","—") or "—",t.get("email",""),t.get("phone",""),
                t.get("created_at",""),t.get("completed_at","") or "—",
                t.get("description",""),t.get("solution","") or "—"]
        for col,val in enumerate(vals,1):
            cell = ws1.cell(row=i,column=col,value=val)
            cell.border = border_all
            cell.alignment = center() if col in (1,5,6,7) else right()
            if col==5:
                cell.font = Font(name="Arial",size=10,bold=True,color=PRI_COLOR.get(pri,MID_BLUE))
                cell.fill = fill(row_bg)
            elif col==6:
                cell.font = Font(name="Arial",size=10,bold=True,color=STATUS_COLOR.get(status,TEXT_DARK))
                cell.fill = fill(STATUS_BG.get(status,row_bg))
            elif col==13:
                # solution column – green tint if has content
                has_sol = val and val != "—"
                cell.font = Font(name="Arial",size=10,bold=has_sol,color=GREEN if has_sol else TEXT_DARK)
                cell.fill = fill(LIGHT_GREEN if has_sol else row_bg)
            else:
                cell.font = cell_font(); cell.fill = fill(row_bg)

    # Summary row
    last = len(tickets)+2
    ws1.row_dimensions[last].height = 28
    ws1.cell(row=last,column=1,value="الإجمالي").font = Font(name="Arial",size=11,bold=True,color=WHITE)
    ws1.cell(row=last,column=1).fill = fill(DARK_BLUE); ws1.cell(row=last,column=1).alignment = center()
    ws1.cell(row=last,column=1).border = border_all
    ws1.cell(row=last,column=2,value=f'=COUNTA(B2:B{last-1})').font = Font(name="Arial",size=11,bold=True,color=WHITE)
    ws1.cell(row=last,column=2).fill = fill(MID_BLUE); ws1.cell(row=last,column=2).alignment = center()
    ws1.cell(row=last,column=2).border = border_all

    # ── Sheet 2: Statistics ──
    ws2 = wb.create_sheet("الإحصائيات"); ws2.sheet_view.rightToLeft = True
    total     = len(tickets)
    completed = sum(1 for t in tickets if t["status"]=="مكتمل")
    in_prog   = sum(1 for t in tickets if t["status"]=="قيد التنفيذ")
    new_cnt   = sum(1 for t in tickets if t["status"]=="جديد")

    def ws2_hdr(row,col,val,bg=DARK_BLUE):
        c2 = ws2.cell(row=row,column=col,value=val)
        c2.font = hdr_font(); c2.fill = fill(bg); c2.alignment = center(); c2.border = border_all
    def ws2_val(row,col,val,bold=False,fg=TEXT_DARK,bg=WHITE):
        c2 = ws2.cell(row=row,column=col,value=val)
        c2.font = Font(name="Arial",size=11,bold=bold,color=fg)
        c2.fill = fill(bg); c2.alignment = center(); c2.border = border_all

    ws2.column_dimensions["A"].width=24; ws2.column_dimensions["B"].width=16; ws2.column_dimensions["C"].width=16; ws2.column_dimensions["D"].width=16

    ws2_hdr(1,1,"ملخص النظام",bg=DARK_BLUE); ws2.merge_cells("A1:D1"); ws2.row_dimensions[1].height=32
    ws2_hdr(2,1,"الإجمالي"); ws2_hdr(2,2,"مكتمل"); ws2_hdr(2,3,"قيد التنفيذ"); ws2_hdr(2,4,"جديد")
    ws2.row_dimensions[2].height=26
    ws2_val(3,1,total,True,TEXT_DARK,LIGHT_BLUE); ws2_val(3,2,completed,True,GREEN,LIGHT_GREEN)
    ws2_val(3,3,in_prog,True,ORANGE,LIGHT_ORANGE); ws2_val(3,4,new_cnt,True,RED,LIGHT_RED)
    ws2.row_dimensions[3].height=30

    # By dept
    ws2.row_dimensions[5].height=28; ws2_hdr(5,1,"البلاغات حسب القسم",bg=GRAY_HDR); ws2.merge_cells("A5:B5")
    ws2_hdr(6,1,"القسم"); ws2_hdr(6,2,"عدد البلاغات")
    dept_map={}
    for t in tickets: dept_map[t["dept"]]=dept_map.get(t["dept"],0)+1
    for i,(dept,cnt) in enumerate(sorted(dept_map.items(),key=lambda x:-x[1]),7):
        bg = GRAY_ROW if i%2==0 else WHITE
        ws2_val(i,1,dept,bg=bg); ws2_val(i,2,cnt,True,MID_BLUE,bg)
        ws2.row_dimensions[i].height=20

    # By tech
    sr=7+len(dept_map)+2; ws2.row_dimensions[sr].height=28
    ws2_hdr(sr,1,"أداء الفنيين",bg=GRAY_HDR); ws2.merge_cells(f"A{sr}:B{sr}")
    ws2_hdr(sr+1,1,"الفني"); ws2_hdr(sr+1,2,"البلاغات المُنجزة")
    tech_map={}
    for t in tickets:
        if t.get("tech"): tech_map[t["tech"]]=tech_map.get(t["tech"],0)+1
    for i,(tech,cnt) in enumerate(sorted(tech_map.items(),key=lambda x:-x[1]),sr+2):
        bg = GRAY_ROW if i%2==0 else WHITE
        ws2_val(i,1,tech,bg=bg); ws2_val(i,2,cnt,True,GREEN,bg); ws2.row_dimensions[i].height=20

    # ── Sheet 3: Notes ──
    ws3 = wb.create_sheet("الملاحظات"); ws3.sheet_view.rightToLeft = True
    write_header(ws3,["رقم البلاغ","موظف","الكاتب","الملاحظة","التاريخ"],[14,18,18,50,18],bg=DARK_BLUE)
    for i,n in enumerate(notes,2):
        bg = GRAY_ROW if i%2==0 else WHITE; ws3.row_dimensions[i].height=22
        for col,val in enumerate([n["ticket_no"],n.get("employee",""),n["author"],n["note"],n["created_at"]],1):
            c3 = ws3.cell(row=i,column=col,value=val)
            c3.font = cell_font(); c3.fill = fill(bg); c3.border = border_all
            c3.alignment = center() if col in (1,3,5) else right()

    # ── Sheet 4: Users ──
    ws4 = wb.create_sheet("المستخدمون"); ws4.sheet_view.rightToLeft = True
    write_header(ws4,["م","اسم المستخدم","الاسم الكامل","الدور","الحالة"],[6,20,24,14,12],bg=DARK_BLUE)
    role_ar = {"admin":"مدير","tech":"فني"}
    for i,u in enumerate([x for x in users if x["active"]],2):
        bg = GRAY_ROW if i%2==0 else WHITE; ws4.row_dimensions[i].height=22
        for col,val in enumerate([i-1,u["username"],u.get("fullname",""),role_ar.get(u["role"],u["role"]),"نشط"],1):
            c4 = ws4.cell(row=i,column=col,value=val)
            c4.font = cell_font(bold=(col==4))
            if col==4:
                c4.font = Font(name="Arial",size=10,bold=True,
                    color=MID_BLUE if u["role"]=="admin" else GREEN)
            c4.fill = fill(bg); c4.border = border_all; c4.alignment = center()

    # ── Sheet 5: Tech Report (تقرير الفنيين) ──
    ws5 = wb.create_sheet("تقرير الفنيين"); ws5.sheet_view.rightToLeft = True

    # Build per-tech statistics
    tech_stats = {}
    for t in tickets:
        tech = t.get("tech","")
        if not tech:
            continue
        if tech not in tech_stats:
            tech_stats[tech] = {
                "total": 0, "completed": 0, "in_progress": 0, "new": 0,
                "resolution_times": [], "tickets": []
            }
        tech_stats[tech]["total"] += 1
        if t["status"] == "مكتمل":
            tech_stats[tech]["completed"] += 1
            if t.get("created_at") and t.get("completed_at"):
                try:
                    t1 = datetime.strptime(t["created_at"][:16],"%Y-%m-%d %H:%M")
                    t2 = datetime.strptime(t["completed_at"][:16],"%Y-%m-%d %H:%M")
                    tech_stats[tech]["resolution_times"].append((t2-t1).total_seconds()/3600)
                except:
                    pass
        elif t["status"] == "قيد التنفيذ":
            tech_stats[tech]["in_progress"] += 1
        else:
            tech_stats[tech]["new"] += 1
        tech_stats[tech]["tickets"].append(t)

    # ── Section 1: Summary header ──
    ws5.column_dimensions["A"].width = 20
    ws5.column_dimensions["B"].width = 14
    ws5.column_dimensions["C"].width = 14
    ws5.column_dimensions["D"].width = 14
    ws5.column_dimensions["E"].width = 16
    ws5.column_dimensions["F"].width = 16
    ws5.column_dimensions["G"].width = 16

    # Title
    title_cell = ws5.cell(row=1, column=1, value="تقرير أداء الفنيين")
    title_cell.font = Font(name="Arial", size=16, bold=True, color=WHITE)
    title_cell.fill = fill(DARK_BLUE)
    title_cell.alignment = center()
    title_cell.border = border_all
    ws5.merge_cells("A1:G1")
    ws5.row_dimensions[1].height = 36

    # Sub-title with date
    date_cell = ws5.cell(row=2, column=1, value=f"تاريخ التصدير: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    date_cell.font = Font(name="Arial", size=10, color=WHITE)
    date_cell.fill = fill(GRAY_HDR)
    date_cell.alignment = center()
    date_cell.border = border_all
    ws5.merge_cells("A2:G2")
    ws5.row_dimensions[2].height = 22

    # ── Section 2: Summary table per tech ──
    summary_headers = ["اسم الفني","إجمالي البلاغات","مكتمل","قيد التنفيذ","جديد","نسبة الإنجاز","متوسط وقت الحل (ساعة)"]
    ws5.row_dimensions[4].height = 28
    for col, hdr in enumerate(summary_headers, 1):
        cell = ws5.cell(row=4, column=col, value=hdr)
        cell.font = hdr_font(sz=11)
        cell.fill = fill(GRAY_HDR)
        cell.alignment = center()
        cell.border = border_all

    row_idx = 5
    for tech, stats in sorted(tech_stats.items(), key=lambda x: -x[1]["completed"]):
        bg = GRAY_ROW if row_idx % 2 == 0 else WHITE
        ws5.row_dimensions[row_idx].height = 22
        total_t = stats["total"]
        comp_t  = stats["completed"]
        rate    = round(comp_t / total_t * 100, 1) if total_t else 0
        times   = stats["resolution_times"]
        avg_t   = round(sum(times)/len(times), 1) if times else "—"

        row_vals = [tech, total_t, comp_t, stats["in_progress"], stats["new"],
                    f"{rate}%", avg_t]
        for col, val in enumerate(row_vals, 1):
            cell = ws5.cell(row=row_idx, column=col, value=val)
            cell.border = border_all
            cell.alignment = center()
            if col == 1:
                cell.font = Font(name="Arial", size=10, bold=True, color=DARK_BLUE)
                cell.fill = fill(LIGHT_BLUE)
            elif col == 3:
                cell.font = Font(name="Arial", size=10, bold=True, color=GREEN)
                cell.fill = fill(LIGHT_GREEN if comp_t > 0 else bg)
            elif col == 4:
                cell.font = Font(name="Arial", size=10, bold=True, color=ORANGE)
                cell.fill = fill(LIGHT_ORANGE if stats["in_progress"] > 0 else bg)
            elif col == 5:
                cell.font = Font(name="Arial", size=10, bold=True, color=RED)
                cell.fill = fill(LIGHT_RED if stats["new"] > 0 else bg)
            elif col == 6:
                rate_val = float(str(val).replace("%","")) if val != "—" else 0
                fg_color = GREEN if rate_val >= 80 else (ORANGE if rate_val >= 50 else RED)
                cell.font = Font(name="Arial", size=10, bold=True, color=fg_color)
                cell.fill = fill(bg)
            else:
                cell.font = cell_font()
                cell.fill = fill(bg)
        row_idx += 1

    # ── Section 3: Detailed tickets per tech ──
    row_idx += 2  # blank separator

    for tech, stats in sorted(tech_stats.items(), key=lambda x: x[0]):
        # Tech name header
        tech_hdr_cell = ws5.cell(row=row_idx, column=1, value=f"بلاغات الفني: {tech}")
        tech_hdr_cell.font = Font(name="Arial", size=12, bold=True, color=WHITE)
        tech_hdr_cell.fill = fill(MID_BLUE)
        tech_hdr_cell.alignment = center()
        tech_hdr_cell.border = border_all
        ws5.merge_cells(f"A{row_idx}:G{row_idx}")
        ws5.row_dimensions[row_idx].height = 26
        row_idx += 1

        # Column headers for tickets
        detail_headers = ["رقم البلاغ","الموظف","القسم","نوع المشكلة","الأولوية","الحالة","الحل المُطبَّق"]
        for col, hdr in enumerate(detail_headers, 1):
            cell = ws5.cell(row=row_idx, column=col, value=hdr)
            cell.font = hdr_font(sz=10)
            cell.fill = fill(GRAY_HDR)
            cell.alignment = center()
            cell.border = border_all
        ws5.row_dimensions[row_idx].height = 22
        row_idx += 1

        # Ticket rows
        for t in stats["tickets"]:
            bg = GRAY_ROW if row_idx % 2 == 0 else WHITE
            ws5.row_dimensions[row_idx].height = 20
            status = t.get("status","جديد")
            pri    = t.get("priority","عادي")
            solution_val = t.get("solution","") or "—"
            row_vals = [
                t["ticket_no"], t["name"], t["dept"], t["issue"],
                pri, status, solution_val
            ]
            for col, val in enumerate(row_vals, 1):
                cell = ws5.cell(row=row_idx, column=col, value=val)
                cell.border = border_all
                cell.alignment = center() if col in (1,5,6) else right()
                if col == 5:
                    cell.font = Font(name="Arial",size=10,bold=True,color=PRI_COLOR.get(pri,MID_BLUE))
                    cell.fill = fill(bg)
                elif col == 6:
                    cell.font = Font(name="Arial",size=10,bold=True,color=STATUS_COLOR.get(status,TEXT_DARK))
                    cell.fill = fill(STATUS_BG.get(status,bg))
                elif col == 7:
                    has_sol = val and val != "—"
                    cell.font = Font(name="Arial",size=10,bold=has_sol,color=GREEN if has_sol else TEXT_DARK)
                    cell.fill = fill(LIGHT_GREEN if has_sol else bg)
                else:
                    cell.font = cell_font()
                    cell.fill = fill(bg)
            row_idx += 1

        row_idx += 2  # blank separator between techs

    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    fname = f"AMP_Support_Export_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"
    return send_file(buf, mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                     as_attachment=True, download_name=fname)


# ══════════════════════════════════════════════
#  EXPORT: PowerPoint
# ══════════════════════════════════════════════
@app.route("/api/export/pptx")
def api_export_pptx():
    if not require_admin(): return jsonify({"error":"unauthorized"}),403
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    conn = get_db(); c = conn.cursor()
    c.execute("SELECT * FROM tickets ORDER BY id"); tickets=[dict(r) for r in c.fetchall()]
    conn.close()

    total     = len(tickets)
    completed = sum(1 for t in tickets if t["status"]=="مكتمل")
    in_prog   = sum(1 for t in tickets if t["status"]=="قيد التنفيذ")
    new_cnt   = sum(1 for t in tickets if t["status"]=="جديد")
    rate      = round(completed/total*100,1) if total else 0

    dept_map={}
    for t in tickets: dept_map[t["dept"]]=dept_map.get(t["dept"],0)+1
    issue_map={}
    for t in tickets: issue_map[t["issue"]]=issue_map.get(t["issue"],0)+1
    tech_map={}
    for t in tickets:
        if t.get("tech"): tech_map[t["tech"]]=tech_map.get(t["tech"],0)+1

    prs = Presentation(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
    blank = prs.slide_layouts[6]  # completely blank

    # ── Colors ──
    C_DARK = RGBColor(0x1E,0x3A,0x5F)
    C_BLUE = RGBColor(0x25,0x63,0xEB)
    C_GREEN= RGBColor(0x05,0x96,0x52)
    C_ORANGE=RGBColor(0xD9,0x77,0x06)
    C_RED  = RGBColor(0xBE,0x12,0x3C)
    C_WHITE= RGBColor(0xFF,0xFF,0xFF)
    C_LIGHT= RGBColor(0xF1,0xF5,0xF9)
    C_TEXT = RGBColor(0x1E,0x29,0x3B)
    C_GRAY = RGBColor(0x64,0x74,0x8B)

    def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, lw=Pt(0)):
        from pptx.util import Emu
        shape = slide.shapes.add_shape(1, l, t, w, h)
        if fill_color:
            shape.fill.solid(); shape.fill.fore_color.rgb=fill_color
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb=line_color; shape.line.width=lw
        else:
            shape.line.fill.background()
        return shape

    def add_text(slide, text, l, t, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
        txBox = slide.shapes.add_textbox(l, t, w, h)
        tf = txBox.text_frame; tf.word_wrap=True
        p = tf.paragraphs[0]; p.alignment=align
        run = p.add_run(); run.text=text
        run.font.size=Pt(size); run.font.bold=bold; run.font.italic=italic
        if color: run.font.color.rgb=color
        return txBox

    def slide_bg(slide, color):
        bg = slide.background; fill = bg.fill
        fill.solid(); fill.fore_color.rgb=color

    # ═══════════════════
    # SLIDE 1 – Cover
    # ═══════════════════
    s1 = prs.slides.add_slide(blank)
    slide_bg(s1, C_DARK)
    add_rect(s1, Inches(0), Inches(0), Inches(0.18), Inches(7.5), C_BLUE)
    add_rect(s1, Inches(0.18), Inches(0), Inches(0.06), Inches(7.5), C_GREEN)
    add_text(s1,"نظام الدعم الفني · AMP مصراتة",Inches(0.5),Inches(1.8),Inches(12),Inches(1.4),
             size=40,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)
    add_text(s1,"تقرير شامل للبلاغات والعمليات",Inches(0.5),Inches(3.1),Inches(12),Inches(0.8),
             size=22,color=RGBColor(0xBF,0xDB,0xFF),align=PP_ALIGN.CENTER)
    add_text(s1,f"تاريخ التصدير: {datetime.now().strftime('%Y-%m-%d')}",Inches(0.5),Inches(4.0),Inches(12),Inches(0.6),
             size=16,color=C_GRAY,align=PP_ALIGN.CENTER)
    add_rect(s1,Inches(0),Inches(6.8),Inches(13.33),Inches(0.7),RGBColor(0x0F,0x17,0x2A))
    add_text(s1,"Al Madina Properties · Misrata",Inches(0.5),Inches(6.82),Inches(12),Inches(0.5),
             size=13,color=C_GRAY,align=PP_ALIGN.CENTER,italic=True)

    # ═══════════════════
    # SLIDE 2 – KPIs
    # ═══════════════════
    s2 = prs.slides.add_slide(blank)
    slide_bg(s2, C_LIGHT)
    add_rect(s2,Inches(0),Inches(0),Inches(13.33),Inches(1.2),C_DARK)
    add_rect(s2,Inches(0),Inches(1.2),Inches(13.33),Inches(0.04),C_BLUE)
    add_text(s2,"ملخص المؤشرات الرئيسية",Inches(0.5),Inches(0.2),Inches(12),Inches(0.8),
             size=28,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)

    cards = [
        (total,   "إجمالي البلاغات",  C_BLUE,  RGBColor(0xDB,0xEA,0xFE)),
        (completed,"مكتملة",           C_GREEN, RGBColor(0xD1,0xFA,0xE5)),
        (in_prog, "قيد التنفيذ",      C_ORANGE,RGBColor(0xFE,0xF3,0xC7)),
        (new_cnt, "جديدة",            C_RED,   RGBColor(0xFF,0xE4,0xE6)),
    ]
    cw = Inches(2.9); gap = Inches(0.2); start_x = Inches(0.5)
    for i,(val,label,color,bg) in enumerate(cards):
        x = start_x + i*(cw+gap)
        box = add_rect(s2, x, Inches(1.5), cw, Inches(2.5), C_WHITE)
        box.line.color.rgb=color; box.line.width=Pt(2)
        add_rect(s2, x, Inches(1.5), cw, Inches(0.12), color)
        add_text(s2, str(val), x, Inches(2.2), cw, Inches(1.0),
                 size=52, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s2, label, x, Inches(3.3), cw, Inches(0.6),
                 size=16, color=C_TEXT, align=PP_ALIGN.CENTER)

    add_rect(s2,Inches(0.5),Inches(4.4),Inches(12.33),Inches(0.7),C_WHITE)
    add_rect(s2,Inches(0.5),Inches(4.4),Inches(12.33*(rate/100)),Inches(0.7),C_GREEN)
    add_text(s2,f"نسبة الإنجاز: {rate}%",Inches(0.5),Inches(5.2),Inches(12),Inches(0.5),
             size=16,bold=True,color=C_TEXT,align=PP_ALIGN.CENTER)

    # ═══════════════════
    # SLIDE 3 – By Dept
    # ═══════════════════
    s3 = prs.slides.add_slide(blank)
    slide_bg(s3, C_LIGHT)
    add_rect(s3,Inches(0),Inches(0),Inches(13.33),Inches(1.2),C_DARK)
    add_rect(s3,Inches(0),Inches(1.2),Inches(13.33),Inches(0.04),C_BLUE)
    add_text(s3,"البلاغات حسب القسم",Inches(0.5),Inches(0.2),Inches(12),Inches(0.8),
             size=28,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)

    sorted_depts = sorted(dept_map.items(),key=lambda x:-x[1])
    max_v = max(v for _,v in sorted_depts) if sorted_depts else 1
    bar_area_w = Inches(8); bar_area_h = Inches(0.45)
    left_label = Inches(0.5); bar_left = Inches(3.5); val_left = Inches(12.0)
    row_h = Inches(0.55); start_y = Inches(1.5)
    for i,(dept,cnt) in enumerate(sorted_depts[:9]):
        y = start_y + i*row_h
        bg = RGBColor(0xF8,0xFA,0xFC) if i%2==0 else C_WHITE
        add_rect(s3, Inches(0.3), y, Inches(12.73), bar_area_h, bg)
        add_text(s3, dept, left_label, y, Inches(2.8), bar_area_h, size=12, color=C_TEXT)
        bar_w = bar_area_w * (cnt/max_v)
        if bar_w > Inches(0.1):
            add_rect(s3, bar_left, y+Inches(0.08), bar_w, Inches(0.3), C_BLUE)
        add_text(s3, str(cnt), val_left, y, Inches(0.8), bar_area_h,
                 size=13,bold=True,color=C_BLUE,align=PP_ALIGN.RIGHT)

    # ═══════════════════
    # SLIDE 4 – By Issue
    # ═══════════════════
    s4 = prs.slides.add_slide(blank)
    slide_bg(s4, C_LIGHT)
    add_rect(s4,Inches(0),Inches(0),Inches(13.33),Inches(1.2),C_DARK)
    add_rect(s4,Inches(0),Inches(1.2),Inches(13.33),Inches(0.04),C_GREEN)
    add_text(s4,"البلاغات حسب نوع المشكلة",Inches(0.5),Inches(0.2),Inches(12),Inches(0.8),
             size=28,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)

    colors6=[C_BLUE,C_GREEN,C_ORANGE,C_RED,RGBColor(0x8B,0x5C,0xF6),RGBColor(0x06,0xB6,0xD4)]
    sorted_issues=sorted(issue_map.items(),key=lambda x:-x[1])
    max_i=max(v for _,v in sorted_issues) if sorted_issues else 1
    for i,(issue,cnt) in enumerate(sorted_issues[:9]):
        y = start_y + i*row_h
        bg = RGBColor(0xF8,0xFA,0xFC) if i%2==0 else C_WHITE
        add_rect(s4, Inches(0.3), y, Inches(12.73), bar_area_h, bg)
        add_text(s4, issue, left_label, y, Inches(2.8), bar_area_h, size=12, color=C_TEXT)
        bar_w = bar_area_w*(cnt/max_i)
        col = colors6[i%len(colors6)]
        if bar_w>Inches(0.1): add_rect(s4, bar_left, y+Inches(0.08), bar_w, Inches(0.3), col)
        add_text(s4,str(cnt),val_left,y,Inches(0.8),bar_area_h,size=13,bold=True,color=col,align=PP_ALIGN.RIGHT)

    # ═══════════════════
    # SLIDE 5 – Tech Performance
    # ═══════════════════
    s5 = prs.slides.add_slide(blank)
    slide_bg(s5, C_LIGHT)
    add_rect(s5,Inches(0),Inches(0),Inches(13.33),Inches(1.2),C_DARK)
    add_rect(s5,Inches(0),Inches(1.2),Inches(13.33),Inches(0.04),C_GREEN)
    add_text(s5,"أداء الفنيين",Inches(0.5),Inches(0.2),Inches(12),Inches(0.8),
             size=28,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)
    sorted_techs=sorted(tech_map.items(),key=lambda x:-x[1])
    max_t=max(v for _,v in sorted_techs) if sorted_techs else 1
    for i,(tech,cnt) in enumerate(sorted_techs):
        y=start_y+i*row_h
        bg=RGBColor(0xF8,0xFA,0xFC) if i%2==0 else C_WHITE
        add_rect(s5,Inches(0.3),y,Inches(12.73),bar_area_h,bg)
        add_text(s5,tech,left_label,y,Inches(2.8),bar_area_h,size=13,color=C_TEXT)
        bw=bar_area_w*(cnt/max_t)
        if bw>Inches(0.1): add_rect(s5,bar_left,y+Inches(0.08),bw,Inches(0.3),C_GREEN)
        add_text(s5,str(cnt),val_left,y,Inches(0.8),bar_area_h,size=14,bold=True,color=C_GREEN,align=PP_ALIGN.RIGHT)

    # ═══════════════════
    # SLIDE 6 – Recent Tickets
    # ═══════════════════
    s6 = prs.slides.add_slide(blank)
    slide_bg(s6, C_LIGHT)
    add_rect(s6,Inches(0),Inches(0),Inches(13.33),Inches(1.2),C_DARK)
    add_rect(s6,Inches(0),Inches(1.2),Inches(13.33),Inches(0.04),C_ORANGE)
    add_text(s6,"آخر البلاغات",Inches(0.5),Inches(0.2),Inches(12),Inches(0.8),
             size=28,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)

    recent=tickets[-10:][::-1]
    cols=[("رقم البلاغ",Inches(1.2)),("الاسم",Inches(2.5)),("المشكلة",Inches(2.5)),("الأولوية",Inches(1.3)),("الحالة",Inches(1.5)),("الفني",Inches(1.8))]
    x_pos=[Inches(0.3)]
    for _,w in cols[:-1]: x_pos.append(x_pos[-1]+w+Inches(0.05))
    for (lbl,w),x in zip(cols,x_pos):
        add_rect(s6,x,Inches(1.35),w,Inches(0.38),C_BLUE)
        add_text(s6,lbl,x,Inches(1.35),w,Inches(0.38),size=11,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)
    STATUS_C={"جديد":C_RED,"قيد التنفيذ":C_ORANGE,"مكتمل":C_GREEN}
    for ri,t in enumerate(recent[:8]):
        ry=Inches(1.35)+Inches(0.38)+ri*Inches(0.45)
        bg=RGBColor(0xF8,0xFA,0xFC) if ri%2==0 else C_WHITE
        row_vals=[t["ticket_no"],t["name"],t["issue"],t.get("priority","عادي"),t["status"],t.get("tech","—") or "—"]
        for ci,((lbl,w),x) in enumerate(zip(cols,x_pos)):
            add_rect(s6,x,ry,w,Inches(0.4),bg)
            color=STATUS_C.get(t["status"],C_TEXT) if ci==4 else C_TEXT
            add_text(s6,str(row_vals[ci]),x,ry,w,Inches(0.4),size=10,color=color,
                     bold=(ci in (0,4)),align=PP_ALIGN.CENTER)

    buf2=io.BytesIO(); prs.save(buf2); buf2.seek(0)
    fname=f"AMP_Support_Report_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    return send_file(buf2,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,download_name=fname)


if __name__ == "__main__":
    init_db()
    app.run(debug=True, port=5000)
